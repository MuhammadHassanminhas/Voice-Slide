"""
VoiceSlide — PPTX Converter
Extracts text and images from a Microsoft PowerPoint (.pptx) file and
converts them into the VoiceSlide JSON format.
"""

import logging
import os
import uuid
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config

logger = logging.getLogger(__name__)


def extract_text_from_shape(shape) -> tuple[str, list[str]]:
    """Extract a heading and list of bullet points from a shape.
    
    Returns:
        (heading_text, list_of_bullet_points)
    """
    heading = ""
    items = []
    
    if not hasattr(shape, "text_frame"):
        return heading, items
        
    for i, paragraph in enumerate(shape.text_frame.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
            
        # First paragraph is often the heading if this is a title placeholder
        # But we rely on the shape name for that usually.
        # Just collect all paragraphs as items for now.
        items.append(text)
        
    return heading, items


def convert_pptx(file_path: str) -> dict[str, Any]:
    """Convert a .pptx file to the VoiceSlide slides.json schema format.

    Args:
        file_path: Absolute path to the .pptx file.

    Returns:
        A dict matching the VoiceSlide slides.json schema.
    """
    logger.info("Converting PPTX file: %s", file_path)
    
    try:
        prs = Presentation(file_path)
    except Exception as exc:
        raise ValueError(f"Failed to parse PPTX file: {exc}") from exc

    title = "Imported Presentation"
    if prs.core_properties.title:
        title = prs.core_properties.title

    slides_data = []
    
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_dict: dict[str, Any] = {
            "id": slide_index,
            "type": "bullets",  # Default type
            "heading": "",
            "items": [],
            "notes": "",
        }

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_dict["notes"] = notes_text

        has_image = False
        text_blocks = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not text:
                    continue
                    
                # Simplistic heuristic: if it's a title shape, it's the heading
                if getattr(shape, "is_placeholder", False) and "title" in shape.name.lower():
                    slide_dict["heading"] = text
                elif getattr(shape, "is_placeholder", False) and "subtitle" in shape.name.lower():
                    if slide_dict["type"] == "bullets":
                        slide_dict["type"] = "title"
                    slide_dict["subheading"] = text
                else:
                    _, items = extract_text_from_shape(shape)
                    text_blocks.extend(items)
                    
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
                # Extract the image
                try:
                    image = shape.image
                    image_ext = image.ext
                    image_bytes = image.blob
                    
                    # Save image to static folder
                    image_filename = f"slide_{slide_index}_{uuid.uuid4().hex[:8]}.{image_ext}"
                    image_path = os.path.join(config.STATIC_DIR, "images", image_filename)
                    
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                        
                    slide_dict["image_url"] = f"/static/images/{image_filename}"
                except Exception as e:
                    logger.warning("Failed to extract image from shape: %s", e)

        # Post-process text blocks into bullet items
        if text_blocks:
             # Just use all text blocks as bullet points if we don't have a specific body text field
             slide_dict["items"] = text_blocks

        # Refine slide type heuristic
        if slide_index == 1 and not slide_dict["items"]:
             slide_dict["type"] = "title"
             # If no heading was found but there are text blocks, use the first as heading
             if not slide_dict["heading"] and text_blocks:
                 slide_dict["heading"] = text_blocks[0]
                 if len(text_blocks) > 1:
                     slide_dict["subheading"] = text_blocks[1]
                 slide_dict["items"] = []
                 
        elif has_image:
             slide_dict["type"] = "image"
             if slide_dict["items"]:
                 slide_dict["caption"] = " ".join(slide_dict["items"])
                 del slide_dict["items"]
                 
        elif len(slide_dict["items"]) == 0 and slide_dict["heading"]:
             slide_dict["type"] = "title"
             
        elif len(slide_dict["items"]) == 1 and len(slide_dict["items"][0]) > 100:
             # Long single block of text is probably a text slide
             slide_dict["type"] = "text"
             slide_dict["body"] = slide_dict["items"][0]
             del slide_dict["items"]

        # Ensure every slide has at least a heading
        if not slide_dict.get("heading"):
            slide_dict["heading"] = f"Slide {slide_index}"

        slides_data.append(slide_dict)

    from datetime import datetime
    
    result = {
        "version": "1.0",
        "title": title,
        "author": prs.core_properties.author or "Unknown",
        "theme": "dark",
        "created_at": datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ"),
        "slides": slides_data
    }
    
    logger.info("Successfully converted %d slides", len(slides_data))
    return result
